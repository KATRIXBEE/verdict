from pptx import Presentation

prs = Presentation(r"C:\Users\ASUS\Downloads\SIH2026-IDEA-Presentation-Format.pptx")

for i in range(1, 6):
    slide = prs.slides[i]
    print(f"\n--- Slide {i+1} Shapes ---")
    for j, shape in enumerate(slide.shapes):
        txt = shape.text_frame.text if shape.has_text_frame else "[NO TEXT]"
        name = shape.name
        is_title = (shape == slide.shapes.title) or "Title" in name
        print(f"  Shape {j}: name='{name}', type={shape.shape_type}, is_title={is_title}, text_preview={repr(txt[:40])}")
