from pptx import Presentation

prs = Presentation(r"C:\Users\ASUS\Downloads\SIH2026-IDEA-Presentation-Format.pptx")

for i, slide in enumerate(prs.slides):
    print(f"\n==================== SLIDE {i+1} ====================")
    for j, shape in enumerate(slide.shapes):
        if shape.has_text_frame:
            print(f"\n--- Shape {j} ({shape.name}) ---")
            for p_idx, p in enumerate(shape.text_frame.paragraphs):
                p_text = "".join(r.text for r in p.runs) if p.runs else p.text
                print(f"  P{p_idx}: {p_text}")
