from pptx import Presentation

prs = Presentation(r"C:\Users\ASUS\Downloads\SIH2026-IDEA-Presentation-Format.pptx")
print(f"Slide Width: {prs.slide_width / 914400:.2f} inches")
print(f"Slide Height: {prs.slide_height / 914400:.2f} inches")
print(f"Total Slides: {len(prs.slides)}")

for i, slide in enumerate(prs.slides):
    print(f"\n=== SLIDE {i+1} ===")
    for j, shape in enumerate(slide.shapes):
        txt = shape.text_frame.text if shape.has_text_frame else "[NO TEXT]"
        txt_preview = repr(txt[:80])
        print(f"  Shape {j} ({shape.name}, type={shape.shape_type}): pos=({shape.left/914400:.2f}, {shape.top/914400:.2f}, w={shape.width/914400:.2f}, h={shape.height/914400:.2f}) | text: {txt_preview}")
