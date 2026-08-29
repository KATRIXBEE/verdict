import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement

template_path = r"C:\Users\ASUS\Downloads\SIH2026-IDEA-Presentation-Format.pptx"
output_path = r"C:\Users\ASUS\OneDrive\Desktop\projects\verdict\VERDICT_SIH2026_Official.pptx"

prs = Presentation(template_path)

# Colors
SIH_BLUE    = (30, 130, 190)
DARK_BLUE   = (0, 60, 130)
NAVY_HEADER = (10, 35, 80)
RED_ACCENT  = (200, 45, 45)
GREEN_ACCENT= (0, 140, 75)
ORANGE_ACC  = (215, 95, 0)
DARK_GREY   = (45, 45, 45)
WHITE       = (255, 255, 255)

# Card Fills
FILL_LIGHT_BLUE   = (238, 246, 255)
FILL_LIGHT_GREEN  = (236, 253, 245)
FILL_LIGHT_YELLOW = (255, 251, 235)
FILL_LIGHT_RED    = (254, 242, 242)
FILL_LIGHT_PURPLE = (245, 243, 255)

def remove_bullets(p):
    pPr = p._p.get_or_add_pPr()
    for child in list(pPr):
        if any(child.tag.endswith(tag) for tag in ['buChar', 'buAutoNum', 'buBlip', 'buClr', 'buFont', 'buSzPct']):
            pPr.remove(child)
    buNone = OxmlElement('a:buNone')
    pPr.append(buNone)

def format_content_textbox(shape, sections_data):
    """
    sections_data: list of tuples:
      (heading_text, list_of_bullet_points, heading_color)
    """
    shape.left = Inches(0.67)
    shape.top = Inches(1.32)
    shape.width = Inches(12.0)
    shape.height = Inches(3.30)
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.clear()
    
    first_p = True
    for sec_idx, (heading, bullets, hcolor) in enumerate(sections_data):
        if heading:
            p = tf.paragraphs[0] if first_p else tf.add_paragraph()
            first_p = False
            remove_bullets(p)
            p.space_before = Pt(5) if sec_idx > 0 else Pt(0)
            p.space_after = Pt(2)
            r = p.add_run()
            r.text = heading
            r.font.size = Pt(10.5)
            r.font.bold = True
            r.font.color.rgb = RGBColor(*hcolor)
            
        for b_text in bullets:
            p = tf.paragraphs[0] if first_p else tf.add_paragraph()
            first_p = False
            remove_bullets(p)
            p.space_before = Pt(0)
            p.space_after = Pt(1.5)
            r = p.add_run()
            r.text = b_text
            r.font.size = Pt(9.0)
            r.font.bold = False
            r.font.color.rgb = RGBColor(*DARK_GREY)

def update_team_oval(slide):
    for shape in slide.shapes:
        if shape.has_text_frame and "Your Team Name" in shape.text_frame.text:
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            remove_bullets(p)
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = "VERDICT"
            r.font.size = Pt(11)
            r.font.bold = True
            r.font.color.rgb = RGBColor(*WHITE)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*DARK_BLUE)
            shape.line.color.rgb = RGBColor(*SIH_BLUE)
            shape.line.width = Pt(1.5)

def add_card(slide, left, top, width, height, title, body_lines, fill_rgb, border_rgb, title_color=DARK_BLUE, title_size=9.5, body_size=8.0):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill_rgb)
    shape.line.color.rgb = RGBColor(*border_rgb)
    shape.line.width = Pt(1.5)
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)
    tf.clear()
    
    p0 = tf.paragraphs[0]
    remove_bullets(p0)
    p0.alignment = PP_ALIGN.CENTER
    p0.space_after = Pt(3)
    r0 = p0.add_run()
    r0.text = title
    r0.font.size = Pt(title_size)
    r0.font.bold = True
    r0.font.color.rgb = RGBColor(*title_color)
    
    for line in body_lines:
        p = tf.add_paragraph()
        remove_bullets(p)
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(1)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(body_size)
        r.font.bold = False
        r.font.color.rgb = RGBColor(*DARK_GREY)
    return shape

def add_arrow_pill(slide, left, top, width=0.4, height=0.3, text="➔", color=SIH_BLUE, font_size=14):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    remove_bullets(p)
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(font_size)
    r.font.bold = True
    r.font.color.rgb = RGBColor(*color)
    return txBox

# ─────────────────────────────────────────────────────────────
# SLIDE 1: TITLE PAGE
# ─────────────────────────────────────────────────────────────
slide1 = prs.slides[0]

# Subtitle shape ("TITLE PAGE") -> Shape 3
sub_shape = slide1.shapes[3]
tf_sub = sub_shape.text_frame
tf_sub.clear()
p = tf_sub.paragraphs[0]
remove_bullets(p)
r = p.add_run()
r.text = "VERDICT"
r.font.size = Pt(32)
r.font.bold = True
r.font.color.rgb = RGBColor(*NAVY_HEADER)

p2 = tf_sub.add_paragraph()
remove_bullets(p2)
r2 = p2.add_run()
r2.text = "India's Politician Accountability & Civic Intelligence Platform"
r2.font.size = Pt(13)
r2.font.bold = True
r2.font.color.rgb = RGBColor(*SIH_BLUE)

p3 = tf_sub.add_paragraph()
remove_bullets(p3)
r3 = p3.add_run()
r3.text = "Rate Your Neta. Know Your Money. Empower Democracy."
r3.font.size = Pt(11)
r3.font.bold = False
r3.font.color.rgb = RGBColor(*DARK_GREY)

# Problem statement details box -> Shape 5
ps_shape = slide1.shapes[5]
tf_ps = ps_shape.text_frame
tf_ps.clear()
rows = [
    ("Problem Statement ID: ", "SIH26102", True, DARK_BLUE),
    ("Problem Statement Title: ", "AI-Powered System to Detect Anomalies, Fraud & Inefficiencies in MPLAD Scheme", True, DARK_BLUE),
    ("Theme: ", "Miscellaneous / Open Innovation", False, DARK_GREY),
    ("PS Category: ", "Software", False, DARK_GREY),
    ("Team ID: ", "SIH2026-TEAM-VERDICT", True, DARK_BLUE),
    ("Team Name: ", "VERDICT", True, DARK_BLUE),
    ("Live Platform: ", "https://verdict.vercel.app", True, GREEN_ACCENT),
]
for i, (label, val, val_bold, val_color) in enumerate(rows):
    p = tf_ps.paragraphs[0] if i == 0 else tf_ps.add_paragraph()
    remove_bullets(p)
    p.space_after = Pt(6)
    r1 = p.add_run()
    r1.text = label
    r1.font.size = Pt(11)
    r1.font.bold = True
    r1.font.color.rgb = RGBColor(*DARK_GREY)
    r2 = p.add_run()
    r2.text = val
    r2.font.size = Pt(11)
    r2.font.bold = val_bold
    r2.font.color.rgb = RGBColor(*val_color)

# ─────────────────────────────────────────────────────────────
# SLIDE 2: IDEA TITLE ("What is VERDICT?")
# ─────────────────────────────────────────────────────────────
slide2 = prs.slides[1]
update_team_oval(slide2)

s2_sections = [
    ("THE CRITICAL CITIZEN PROBLEM IN INDIA TODAY:", [
        "  • Candidate criminal cases & wealth records are locked inside 400-page scanned affidavits and complex legal terminology.",
        "  • Over ₹4.83 Lakh Crore in government fund misuse is buried inside technical CAG audit reports that voters never see.",
        "  • Citizens have no single place before entering the voting booth to check: 'Is my MP performing well?'",
    ], RED_ACCENT),
    ("HOW VERDICT SOLVES THIS (ALL-IN-ONE CITIZEN INTELLIGENCE PORTAL):", [
        "  • 0.0 to 10.0 VERDICT Score: A credit-score-like rating for all 543 Lok Sabha MPs based on verified public records.",
        "  • Plain-Language Criminal Translator: Automatic translation of IPC/BNS sections into 4 severity tiers and layman English.",
        "  • Ground Truth Money Trail: Interactive tracking of CAG audit findings, fund diversions & international cost benchmarks.",
        "  • Neta Face-Off: Side-by-side comparative matrix evaluating any two MPs on attendance, debates, wealth, and cases.",
    ], GREEN_ACCENT),
]
format_content_textbox(slide2.shapes[2], s2_sections)

# Slide 2 Bottom Diagram
y_s2 = 4.75
w_card_s2 = 2.45
h_card_s2 = 1.65

add_card(slide2, 0.67, y_s2, w_card_s2, h_card_s2,
         "1. Citizen Query",
         ["Voter seeks real record", "of candidate before elections", "without political spin"],
         FILL_LIGHT_RED, RED_ACCENT, title_color=RED_ACCENT)

add_arrow_pill(slide2, 0.67 + w_card_s2 + 0.12, y_s2 + 0.65)

add_card(slide2, 3.82, y_s2, w_card_s2, h_card_s2,
         "2. VERDICT Engine",
         ["Ingests ECI affidavits,", "eCourts live dockets,", "Sansad & CAG audit files"],
         FILL_LIGHT_BLUE, SIH_BLUE, title_color=DARK_BLUE)

add_arrow_pill(slide2, 3.82 + w_card_s2 + 0.12, y_s2 + 0.65)

add_card(slide2, 6.97, y_s2, w_card_s2, h_card_s2,
         "3. 0–10 Trust Score",
         ["Instant plain-language score:", "Attendance, criminal charges,", "and wealth growth alert"],
         FILL_LIGHT_YELLOW, ORANGE_ACC, title_color=ORANGE_ACC)

add_arrow_pill(slide2, 6.97 + w_card_s2 + 0.12, y_s2 + 0.65)

add_card(slide2, 10.12, y_s2, w_card_s2, h_card_s2,
         "4. Informed Ballot",
         ["Democracy empowered:", "Citizen votes based on", "hard verified facts"],
         FILL_LIGHT_GREEN, GREEN_ACCENT, title_color=GREEN_ACCENT)

# ─────────────────────────────────────────────────────────────
# SLIDE 3: TECHNICAL APPROACH
# ─────────────────────────────────────────────────────────────
slide3 = prs.slides[2]
update_team_oval(slide3)

s3_sections = [
    ("TECH STACK & ARCHITECTURE (Simple, Modern & Production-Ready):", [
        "  • Web & Mobile UI: Next.js 15 (App Router) + React 19 + Tailwind CSS — responsive, fast PWA with 44px touch targets.",
        "  • Secure Database: PostgreSQL via Supabase with Row-Level Security (RLS), atomic RPCs, and zero client-side secrets.",
        "  • Automated Ingestion: Python pipelines (BeautifulSoup, pandas, requests) scraping ECI, eCourts & Sansad.in.",
        "  • Data Visualizations: Recharts for interactive asset growth trends, budget flows, and CAG category donut charts.",
        "  • Security & SRE: SSRF proxy validation, sliding-window rate limiting, Sentry telemetry, and 0-emoji Lucide icons.",
    ], DARK_BLUE),
    ("MATHEMATICAL SCORING FORMULATION (0.0 to 10.0 Scale):", [
        "  • Base Score: 5.0 / 10.0 (all elected MPs start with an equal neutral baseline).",
        "  • Positive Vectors: +2.0 (Parliament Attendance ≥90%) | +1.0 (Clean Legal Record) | +0.5 (Verified Degree & Loyalty).",
        "  • Penal Deductions: Up to -4.0 for Severe Criminal Prosecutions (Murder, Corruption, Rape) | -2.0 for >500% Asset Outliers.",
    ], DARK_BLUE),
]
format_content_textbox(slide3.shapes[2], s3_sections)

# Slide 3 Bottom Diagram
y_s3 = 4.75
w_s3_box = 1.95
h_s3_box = 1.65

sources = [
    ("ECI Form 26", ["Assets, education,", "affidavits & filings"], (0, 80, 160)),
    ("eCourts NJDG", ["Live court case status,", "FIRs & IPC sections"], (0, 120, 90)),
    ("Sansad / PRS", ["Parliament attendance,", "debates & questions"], (180, 90, 0)),
    ("CAG Audits", ["Official performance", "audit reports & data"], (180, 40, 40)),
]

for idx, (stitle, sbody, sclr) in enumerate(sources):
    x_box = 0.67 + idx * (w_s3_box + 0.12)
    add_card(slide3, x_box, y_s3, w_s3_box, h_s3_box,
             stitle, sbody, FILL_LIGHT_BLUE, sclr, title_color=sclr, title_size=9.5, body_size=8.0)

add_arrow_pill(slide3, 9.0, y_s3 + 0.65)

add_card(slide3, 9.45, y_s3, 3.2, h_s3_box,
         "VERDICT Platform Engine",
         ["• Normalization & Score Calculator", "• IPC Plain-Language Translator", "• Static HTML Pre-render (593 pages)", "• Citizen Dashboard & API Gateway"],
         FILL_LIGHT_GREEN, GREEN_ACCENT, title_color=GREEN_ACCENT, title_size=10, body_size=8.0)

# ─────────────────────────────────────────────────────────────
# SLIDE 4: FEASIBILITY AND VIABILITY
# ─────────────────────────────────────────────────────────────
slide4 = prs.slides[3]
update_team_oval(slide4)

s4_sections = [
    ("IS VERDICT FEASIBLE? YES — PRODUCTION SYSTEM IS ALREADY LIVE & TESTED:", [
        "  • Live Production Build: 593 static routes pre-rendered on Vercel at https://verdict.vercel.app with 0 build errors.",
        "  • 563 Real Politicians Ingested: Complete profiles for 543 Lok Sabha MPs with verified photos, portfolios & asset history.",
        "  • Rigorous SRE Testing: 49/49 automated unit & integration tests passing; P95 response latency <45ms under 100 VUs.",
        "  • Low Cost & Scalable: Entire serverless architecture operates under ₹1,500/month on cloud edge tiers.",
    ], DARK_BLUE),
    ("RISK MITIGATION STRATEGIES (Overcoming Core Real-World Challenges):", [
        "  • Data Scarcity Risk: Solved via automated multi-source scrapers with robust fallback mock caches for offline resilience.",
        "  • Legal / Defamation Risk: Solved by strictly presenting verbatim official records (CAG reports, Supreme Court orders, ECI Form 26).",
        "  • Bot Rating Spam Risk: Solved via server-side session checks, IP-based sliding window rate limits, and Zod input validation.",
    ], DARK_BLUE),
]
format_content_textbox(slide4.shapes[2], s4_sections)

# Slide 4 Bottom Diagram
y_s4 = 4.75
w_s4 = 2.82
roadmap = [
    ("PHASE 1 (LIVE TODAY)", ["v1.0 Operational", "• 543 Lok Sabha MPs", "• 10 CAG Scam Dossiers", "• 593 Static Pages Built"], GREEN_ACCENT, FILL_LIGHT_GREEN),
    ("PHASE 2 (6 MONTHS)", ["v2.0 State Expansion", "• 4,123 State MLAs", "• 28 State Vidhan Sabhas", "• State Budget Audit Tracker"], SIH_BLUE, FILL_LIGHT_BLUE),
    ("PHASE 3 (12 MONTHS)", ["v3.0 Citizen Identity", "• 12 Indian Languages", "• DigiLocker 1-Citizen-1-Vote", "• Constituency WhatsApp Alerts"], ORANGE_ACC, FILL_LIGHT_YELLOW),
    ("PHASE 4 (24 MONTHS)", ["v4.0 Ground Truth AI", "• Satellite Project Audits", "• Live Sansad TV Speech NLP", "• Blockchain Proof Ledger"], (130, 40, 130), FILL_LIGHT_PURPLE),
]

for idx, (rtitle, rbody, rclr, rfill) in enumerate(roadmap):
    x_pos = 0.67 + idx * (w_s4 + 0.22)
    add_card(slide4, x_pos, y_s4, w_s4, 1.65,
             rtitle, rbody, rfill, rclr, title_color=rclr, title_size=9.5, body_size=8.0)
    if idx < 3:
        add_arrow_pill(slide4, x_pos + w_s4 + 0.04, y_s4 + 0.65)

# ─────────────────────────────────────────────────────────────
# SLIDE 5: IMPACT AND BENEFITS
# ─────────────────────────────────────────────────────────────
slide5 = prs.slides[4]
update_team_oval(slide5)

s5_sections = [
    ("MEASURABLE IMPACT & BENEFICIARIES ACROSS INDIAN SOCIETY:", [
        "  • 96 Crore Registered Voters: Enables instant 10-second candidate audits before voting — ending democratic information asymmetry.",
        "  • Taxpayers & Citizens: Demystifies the ₹47.94 Lakh Crore Union Budget, tracking exactly where funds flowed or stalled.",
        "  • Grassroots RTI Activists: 1-click Section 6(1) RTI Form 'A' application generator targeting local Public Information Officers.",
        "  • Investigative Journalists: Centralized, source-linked repository of MP asset histories, party transitions, and court dockets.",
    ], DARK_BLUE),
    ("SYSTEMIC DEMOCRATIC & NATIONAL BENEFITS:", [
        "  • Disincentivizes Crime in Politics: Publicly exposes severe criminal charges and tracks disproportional wealth growth.",
        "  • Supports 'Digital India' & Open Governance: Converts fragmented, hard-to-read government PDFs into structured citizen insights.",
        "  • Promotes Issue-Based Politics: Moves electoral discourse away from divisive rhetoric toward verifiable performance metrics.",
    ], DARK_BLUE),
]
format_content_textbox(slide5.shapes[2], s5_sections)

# Slide 5 Bottom Diagram
y_s5 = 4.75
w_s5 = 2.82
kpis = [
    ("96 CRORE VOTERS", ["Empowered before elections", "with instant 0–10 scorecards", "and plain-language records"], (0, 80, 160), FILL_LIGHT_BLUE),
    ("543 MPs PROFILED", ["100% 18th Lok Sabha", "with attendance, wealth history", "and court case dockets"], (0, 130, 60), FILL_LIGHT_GREEN),
    ("₹4.83 LAKH CRORE", ["Public capital audited", "across 10 verified CAG cases", "with loss translation math"], (200, 45, 45), FILL_LIGHT_RED),
    ("100+ GLOBAL INDICES", ["India's rank vs world", "across 14 domains with", "20-yr Rupee currency ledger"], (180, 90, 0), FILL_LIGHT_YELLOW),
]

for idx, (ktitle, kbody, kclr, kfill) in enumerate(kpis):
    x_pos = 0.67 + idx * (w_s5 + 0.22)
    add_card(slide5, x_pos, y_s5, w_s5, 1.65,
             ktitle, kbody, kfill, kclr, title_color=kclr, title_size=10, body_size=8.0)

# ─────────────────────────────────────────────────────────────
# SLIDE 6: RESEARCH AND REFERENCES
# ─────────────────────────────────────────────────────────────
slide6 = prs.slides[5]
update_team_oval(slide6)

s6_sections = [
    ("OFFICIAL GOVERNMENT DATA SOURCES & STATUTORY REFERENCES:", [
        "  • Election Commission of India (ECI): Form 26 candidate affidavits, assets & education declarations (affidavit.eci.gov.in)",
        "  • National Judicial Data Grid (NJDG) / eCourts: Live district court case status, FIRs & IPC sections (ecourts.gov.in)",
        "  • Sansad.in & PRS Legislative Research: Official Lok Sabha MP attendance transcripts, debate counts & questions asked",
        "  • CAG Audit Reports (cag.gov.in): Report 16/2023 (Dwarka Expressway), Report 7/2023 (PM-JAY), Report 34/2017 (Clean Energy), Minor Head 800",
        "  • Supreme Court of India: Landmark orders in WP(C) 318/2006 (BOCW Cess Fund) & WP(C) 202/1995 (CAMPA Afforestation)",
        "  • Global Governance Indices: World Bank, IMF, UNDP Human Development Index, RSF Press Freedom, WEF Gender Gap, WHO Health",
    ], DARK_BLUE),
    ("LIVE PLATFORM REPOSITORIES & DOCUMENTATION:", [
        "  • Live Production URL: https://verdict.vercel.app  |  Source Code: https://github.com/KATRIXBEE/verdict",
        "  • Interactive REST API Docs: https://verdict.vercel.app/api-docs (OpenAPI / Swagger UI Documentation)",
    ], DARK_BLUE),
]
format_content_textbox(slide6.shapes[2], s6_sections)

# Slide 6 Bottom Badges
y_s6 = 4.85
w_s6 = 2.2
badges = [
    ("ECI Verified", "Form 26 Affidavits", (0, 80, 160)),
    ("eCourts Mapped", "Live Case Dockets", (0, 130, 60)),
    ("CAG Audited", "Official CAG Reports", (200, 45, 45)),
    ("Sansad Certified", "Sansad.in Transcripts", (180, 90, 0)),
    ("Open Source", "Auditable Codebase", (100, 40, 140)),
]

for idx, (btitle, bsub, bclr) in enumerate(badges):
    x_pos = 0.67 + idx * (w_s6 + 0.25)
    add_card(slide6, x_pos, y_s6, w_s6, 1.45,
             btitle, [bsub, "Verified Public Record", "Article 19(1)(a)"], WHITE, bclr, title_color=bclr, title_size=9.5, body_size=7.5)

# ─────────────────────────────────────────────────────────────
# DELETE SLIDE 7 (Important Instructions)
# ─────────────────────────────────────────────────────────────
slide_list = prs.slides._sldIdLst
slide_ids = list(slide_list)
if len(slide_ids) >= 7:
    slide_id_elem = slide_ids[6]
    rId_to_remove = slide_id_elem.get('r:id')
    slide_list.remove(slide_id_elem)
    print(f"Successfully deleted slide 7.")

print(f"Final presentation has {len(prs.slides)} slides.")

# Save files
prs.save(output_path)
prs.save(r"C:\Users\ASUS\OneDrive\Desktop\VERDICT_SIH2026_Official.pptx")
prs.save(r"C:\Users\ASUS\OneDrive\Desktop\projects\VERDICT_SIH2026_Official.pptx")
os.makedirs(r"C:\mnt\user-data\outputs", exist_ok=True)
prs.save(r"C:\mnt\user-data\outputs\VERDICT_SIH2026_Official.pptx")
print("Saved PPTX to all destinations!")
